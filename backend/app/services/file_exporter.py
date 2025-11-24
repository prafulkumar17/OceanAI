from docx import Document as DocxDocument
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches
from typing import Dict, Any
import io
from app.models.project import DocumentType

class FileExporter:
    """Export generated content to actual files"""
    
    @staticmethod
    def export_to_file(content_data: Dict[str, Any], document_type: DocumentType) -> bytes:
        """Export content to file bytes"""
        if document_type == DocumentType.DOCX:
            return FileExporter._export_to_docx(content_data)
        elif document_type == DocumentType.PPTX:
            return FileExporter._export_to_pptx(content_data)
        else:
            raise ValueError(f"Unsupported document type: {document_type}")
    
    @staticmethod
    def _export_to_docx(content_data: Dict[str, Any]) -> bytes:
        """Create a Word document from content"""
        doc = DocxDocument()
        
        # Add title
        if "title" in content_data:
            title = doc.add_heading(content_data["title"], 0)
        
        # Add sections
        sections = content_data.get("sections", [])
        for section in sections:
            # Add section heading
            doc.add_heading(section.get("title", ""), level=1)
            
            # Add paragraphs
            paragraphs = section.get("content", [])
            for para_text in paragraphs:
                para = doc.add_paragraph(para_text)
                para_format = para.paragraph_format
                para_format.space_after = Pt(12)
        
        # Save to bytes
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        return file_stream.read()
    
    @staticmethod
    def _export_to_pptx(content_data: Dict[str, Any]) -> bytes:
        """Create a PowerPoint presentation using Google Slides API"""
        try:
            from app.services.google_slides_service import GoogleSlidesService
            import os
            
            # Get template ID from environment
            template_id = os.getenv("GOOGLE_SLIDES_TEMPLATE_ID")
            if not template_id:
                # Fallback to old method if no template ID
                return FileExporter._export_to_pptx_legacy(content_data)
                
            service = GoogleSlidesService()
            
            # Create new presentation from template
            title = content_data.get("title", "Generated Presentation")
            presentation_id = service.create_presentation_from_template(title, template_id)
            
            try:
                # Prepare replacements
                replacements = {}
                slides = content_data.get("slides", [])
                
                # We assume the template has placeholders like {{title_1}}, {{body_1}}, etc.
                # Or we can duplicate slides. For now, let's assume a simple replacement strategy
                # where we map content to generic placeholders if they exist.
                
                # BETTER STRATEGY:
                # 1. Get the slides of the new presentation
                # 2. Duplicate the "Content" slide (assuming index 1) for each generated slide
                # 3. Replace text in each duplicated slide
                
                # For this MVP, let's just do a simple replacement if the user provided a template
                # designed for this specific content structure.
                
                # Let's assume the template has 1 title slide and 1 content slide.
                # We will duplicate the content slide (index 1) N times.
                
                presentation_slides = service.get_presentation_slides(presentation_id)
                if len(presentation_slides) < 2:
                    raise ValueError("Template must have at least 2 slides")
                    
                title_slide_id = presentation_slides[0]['objectId']
                
                # Find the content slide template (the one with {{SLIDE_CONTENT}})
                content_slide_id = None
                
                # We need to inspect slides to find the placeholder.
                # Since we can't easily inspect text via API without complex calls,
                # we will try to find it by checking if we can replace text on it?
                # No, that's destructive.
                
                # ALTERNATIVE: Just assume it's the one with {{SLIDE_CONTENT}} if the user followed instructions.
                # But we don't know which index it is.
                
                # Let's use a heuristic:
                # If there are > 1 slides, use the LAST slide as the content template?
                # Or use index 1 if available?
                
                # User's template has it at index 3.
                # Let's try to find it by "Layout" name if possible? No.
                
                # Let's iterate and check text?
                # The service.get_presentation_slides returns the slide objects.
                # We can check the text elements in the response!
                
                # Find the content slide template
                # Strategy: Look for a slide that has {{SLIDE_CONTENT}}.
                # Preference: A slide that has BOTH {{SLIDE_TITLE}} and {{SLIDE_CONTENT}}.
                
                content_slide_id = None
                best_score = 0
                
                for slide in presentation_slides:
                    score = 0
                    has_content = False
                    has_title = False
                    
                    # Check text elements
                    for element in slide.get('pageElements', []):
                        shape = element.get('shape')
                        if shape and shape.get('text'):
                            for te in shape['text']['textElements']:
                                if 'textRun' in te:
                                    content = te['textRun']['content']
                                    if '{{SLIDE_CONTENT}}' in content:
                                        has_content = True
                                    if '{{SLIDE_TITLE}}' in content:
                                        has_title = True
                    
                    if has_content:
                        score += 1
                    if has_title:
                        score += 1
                        
                    if score > best_score:
                        best_score = score
                        content_slide_id = slide['objectId']
                        
                    # If we found a perfect match (score 2), stop.
                    if best_score == 2:
                        break
                
                # Fallback if not found: Use index 1 (or 3 if we want to be specific to this user, but let's stick to 1 as default)
                if not content_slide_id:
                     # If user has > 3 slides, maybe try 3?
                     if len(presentation_slides) > 3:
                         content_slide_id = presentation_slides[3]['objectId']
                     else:
                         content_slide_id = presentation_slides[1]['objectId']
                
                # Find the content slide index
                content_slide_index = 0
                for i, slide in enumerate(presentation_slides):
                    if slide['objectId'] == content_slide_id:
                        content_slide_index = i
                        break
                
                # Update Title Slide
                # We use a specific unique placeholder for the main title
                service.replace_text(presentation_id, {
                    "{{MAIN_TITLE}}": title,
                    "{{SUBTITLE}}": "Generated by OceanAI"
                }, slide_ids=[title_slide_id])
                
                # Create slides and update content
                # Strategy: Duplicate the "clean" template slide for EACH generated slide.
                # Insert them after the template slide.
                # Finally, delete the original template slide.
                
                for i, slide_data in enumerate(slides):
                    # Duplicate content slide
                    # Insert at: content_slide_index + 1 + i
                    # e.g. if template is at 1.
                    # i=0 -> insert at 2.
                    # i=1 -> insert at 3.
                    insertion_index = content_slide_index + 1 + i
                    new_slide_id = service.duplicate_slide(presentation_id, content_slide_id, insertion_index=insertion_index)
                    
                    # Prepare content text (join bullets)
                    bullets = slide_data.get("bullets", [])
                    content_text = "\n".join(bullets) if bullets else ""
                    
                    # Replace text ONLY on this new slide
                    service.replace_text(presentation_id, {
                        "{{SLIDE_TITLE}}": slide_data.get("title", ""),
                        "{{SLIDE_CONTENT}}": content_text
                    }, slide_ids=[new_slide_id])

                # Remove the original template content slide
                service.delete_slide(presentation_id, content_slide_id)

            except Exception as e:
                # Cleanup on error
                service.delete_file(presentation_id)
                raise e
                
            # Export
            pptx_bytes = service.export_presentation(presentation_id)
            
            # Cleanup (delete the temp file from Drive)
            service.delete_file(presentation_id)
            
            return pptx_bytes
            
        except Exception as e:
            print(f"Google Slides API failed: {e}. Falling back to legacy.")
            return FileExporter._export_to_pptx_legacy(content_data)

    @staticmethod
    def _export_to_pptx_legacy(content_data: Dict[str, Any]) -> bytes:
        """Legacy method: Create a PowerPoint presentation from content using python-pptx"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add slides
        slides = content_data.get("slides", [])
        for slide_data in slides:
            # Add slide
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = slide_data.get("title", "")
            
            # Add content
            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            tf.word_wrap = True
            tf.clear()  # Clear default text
            
            bullets = slide_data.get("bullets", [])
            if isinstance(bullets, list) and bullets:
                for i, bullet_text in enumerate(bullets):
                    if bullet_text:
                        p = tf.add_paragraph()
                        p.text = bullet_text
                        p.level = 0
        
        # Save to bytes
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        return file_stream.read()

